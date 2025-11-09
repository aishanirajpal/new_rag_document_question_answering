import cohere
import fitz
from pinecone import Pinecone, ServerlessSpec
import io
import os
import docx
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.documents import Document

class VectorStore:
    def __init__(self, doc_files: list, cohere_api_key: str = None, pinecone_api_key: str = None):
        self.doc_files = doc_files
        self.cohere_api_key = cohere_api_key
        self.pinecone_api_key = pinecone_api_key
        
        if self.cohere_api_key and self.pinecone_api_key:
            self.co = cohere.Client(cohere_api_key)
            self.use_hugging_face = False
        else:
            self.embedding_model = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
            self.use_hugging_face = True

        self.chunks = []
        self.embeddings = []
        self.retrieve_top_k = 10
        self.rerank_top_k = 3
        self.doc_texts = []
        self.load_documents()
        self.split_text()
        self.embed_chunks()
        self.index_chunks()

    def load_documents(self):
        for doc_file in self.doc_files:
            # First, rewind the file stream to the beginning.
            doc_file.seek(0)
            # Second, read the entire file content into a new BytesIO object.
            # This creates a standardized, reliable in-memory file for all parsers.
            file_stream = io.BytesIO(doc_file.read())

            file_ext = os.path.splitext(doc_file.name)[1].lower()
            
            text_extractor = getattr(self, f"extract_text_from_{file_ext[1:]}", None)
            
            if text_extractor:
                text = text_extractor(file_stream, doc_file.name)
                if text and text.get("text", "").strip(): # Ensure extracted text is not empty
                    self.doc_texts.append(text)

    def extract_text_from_pdf(self, file_stream: io.BytesIO, file_name: str) -> dict:
        text = ""
        # fitz requires the raw bytes, so we use getvalue()
        with fitz.open(stream=file_stream.getvalue(), filetype="pdf") as pdf:
            for page in pdf:
                text += page.get_text()
        return {"text": text, "source": file_name}

    def extract_text_from_txt(self, file_stream: io.BytesIO, file_name: str) -> dict:
        """Reads a text file, attempting to decode with utf-8 and falling back to latin-1."""
        try:
            text = file_stream.read().decode('utf-8')
        except UnicodeDecodeError:
            # Rewind the stream after a failed read and try another common encoding
            file_stream.seek(0)
            try:
                text = file_stream.read().decode('latin-1')
            except Exception:
                text = "" # Return empty string if all decoding attempts fail
        return {"text": text, "source": file_name}

    def extract_text_from_docx(self, file_stream: io.BytesIO, file_name: str) -> dict:
        """Extracts text from paragraphs and tables in a .docx file."""
        document = docx.Document(file_stream)
        text_parts = []
        for para in document.paragraphs:
            text_parts.append(para.text)
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    text_parts.append(cell.text)
        return {"text": "\n".join(text_parts), "source": file_name}

    def extract_text_from_pptx(self, file_stream: io.BytesIO, file_name: str) -> dict:
        """Extracts text from text frames and tables in a .pptx file."""
        presentation = Presentation(file_stream)
        text_parts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_parts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text_parts.append(cell.text_frame.text)
        return {"text": "\n".join(text_parts), "source": file_name}

    def split_text(self, chunk_size=1000):
        for doc in self.doc_texts:
            sentences = doc["text"].split(". ")
            current_chunk = ""
            for sentence in sentences:
                if len(current_chunk) + len(sentence) < chunk_size:
                    current_chunk += sentence + ". "
                else:
                    self.chunks.append({"text": current_chunk, "source": doc["source"]})
                    current_chunk = sentence + ". "
            if current_chunk:
                self.chunks.append({"text": current_chunk, "source": doc["source"]})

    def embed_chunks(self, batch_size=90):
        if self.use_hugging_face:
            texts = [chunk['text'] for chunk in self.chunks]
            self.embeddings = self.embedding_model.embed_documents(texts)
        else:
            total_chunks = len(self.chunks)
            for i in range(0, total_chunks, batch_size):
                batch_texts = [chunk['text'] for chunk in self.chunks[i:min(i + batch_size, total_chunks)]]
                batch_embeddings = self.co.embed(
                    texts=batch_texts, input_type="search_document", model="embed-english-v3.0"
                ).embeddings
                self.embeddings.extend(batch_embeddings)

    def index_chunks(self):
        """
        Indexes the embedded chunks.
        """
        if not self.embeddings:
            return

        if self.use_hugging_face:
            documents = []
            for i, chunk in enumerate(self.chunks):
                doc = Document(page_content=chunk['text'], metadata={'source': chunk['source']})
                documents.append(doc)
            
            self.index = FAISS.from_documents(documents, self.embedding_model)
        else:
            pc = Pinecone(
                api_key=self.pinecone_api_key
            )

            index_name = 'rag-qa-bot'
            if index_name not in pc.list_indexes().names():
                pc.create_index(
                    name=index_name,
                    dimension=len(self.embeddings[0]),
                    metric='cosine',
                    spec=ServerlessSpec(
                        cloud='aws',
                        region='us-east-1'
                    )
                )
            self.index = pc.Index(index_name)

            # Clear the index before upserting new data
            if self.index.describe_index_stats()['total_vector_count'] > 0:
                self.index.delete(delete_all=True)

            chunks_metadata = [{'text': chunk['text'], 'source': chunk['source']} for chunk in self.chunks]
            ids = [str(i) for i in range(len(self.chunks))]
            self.index.upsert(vectors=zip(ids, self.embeddings, chunks_metadata))

    def retrieve(self, query: str) -> list:
        if self.use_hugging_face:
            results = self.index.similarity_search(query, k=self.retrieve_top_k)
            return [{'text': doc.page_content, 'source': doc.metadata['source']} for doc in results]
        else:
            query_emb = self.co.embed(
                texts=[query], model="embed-english-v3.0", input_type="search_query"
            ).embeddings
            res = self.index.query(vector=query_emb, top_k=self.retrieve_top_k, include_metadata=True)
            
            docs_to_rerank = []
            for match in res['matches']:
                docs_to_rerank.append({'text': match['metadata']['text'], 'source': match['metadata']['source']})

            rerank_results = self.co.rerank(
                query=query,
                documents=[doc['text'] for doc in docs_to_rerank],
                top_n=self.rerank_top_k,
                model="rerank-english-v3.0"
            )
            
            final_docs = []
            for result in rerank_results.results:
                original_doc = docs_to_rerank[result.index]
                final_docs.append({'text': original_doc['text'], 'source': original_doc['source']})
                
            return final_docs

    def get_document_text(self, file_name: str) -> str:
        for doc in self.doc_texts:
            if doc['source'] == file_name:
                return doc['text']
        return None
